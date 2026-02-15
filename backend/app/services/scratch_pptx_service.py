from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

_BULLET_PREFIX = re.compile(r"^\s*(?:\u2022|-|\*)\s*")
_THEMES_DIR = Path(__file__).with_name("themes")
_THEME_COLORS = {
    "default": {
        "bg": RGBColor(236, 243, 253),
        "title": RGBColor(15, 23, 42),
        "body": RGBColor(30, 41, 59),
        "accent": RGBColor(29, 78, 216),
        "accent_soft": RGBColor(191, 219, 254),
        "card_bg": RGBColor(255, 255, 255),
        "card_border": RGBColor(147, 197, 253),
    },
    "dark": {
        "bg": RGBColor(15, 23, 42),
        "title": RGBColor(248, 250, 252),
        "body": RGBColor(226, 232, 240),
        "accent": RGBColor(56, 189, 248),
        "accent_soft": RGBColor(30, 41, 59),
        "card_bg": RGBColor(30, 41, 59),
        "card_border": RGBColor(51, 65, 85),
    },
    "corporate": {
        "bg": RGBColor(238, 245, 251),
        "title": RGBColor(15, 23, 42),
        "body": RGBColor(30, 64, 96),
        "accent": RGBColor(14, 116, 144),
        "accent_soft": RGBColor(191, 219, 254),
        "card_bg": RGBColor(255, 255, 255),
        "card_border": RGBColor(165, 204, 236),
    },
}


def _normalize_lines(text: str) -> list[str]:
    lines: list[str] = []
    for raw in (text or "").splitlines():
        line = raw.strip()
        if not line:
            continue
        line = _BULLET_PREFIX.sub("", line).strip()
        if line:
            lines.append(line)
    return lines


def _presentation_for_theme(theme: str) -> Presentation:
    theme_name = (theme or "default").strip().lower()
    theme_path = _THEMES_DIR / f"{theme_name}.pptx"
    if theme_path.exists():
        return Presentation(str(theme_path))
    return Presentation()


def _find_layout(prs: Presentation, keywords: tuple[str, ...], fallback_index: int) -> int:
    for idx, layout in enumerate(prs.slide_layouts):
        name = str(getattr(layout, "name", "")).lower()
        if all(word in name for word in keywords):
            return idx
    return min(fallback_index, len(prs.slide_layouts) - 1)


def _pick_layout_index(prs: Presentation, archetype: str, slide_index: int) -> int:
    if slide_index == 0:
        return _find_layout(prs, ("title", "slide"), 0)
    if archetype == "section_break":
        return _find_layout(prs, ("section", "header"), 2)
    if archetype == "comparison":
        return _find_layout(prs, ("two", "content"), 3)
    if archetype == "timeline":
        return _find_layout(prs, ("title", "and", "content"), 1)
    if archetype == "table_data":
        return _find_layout(prs, ("content", "with", "caption"), 7)
    return _find_layout(prs, ("title", "and", "content"), 1)


def _adjust_layout_for_content(prs: Presentation, slots: dict[str, str]) -> int | None:
    bullet_count = len([row for row in _normalize_lines(slots.get("BULLET_1", "")) if row])
    body_length = len((slots.get("BODY_1", "") or "").strip())

    if bullet_count >= 5 or body_length > 300:
        return _find_layout(prs, ("title", "and", "content"), 1)
    if bullet_count <= 2 and body_length < 100:
        return _find_layout(prs, ("title", "and", "content"), 1)
    return None


def _shape_name(shape) -> str:
    return str(getattr(shape, "name", "")).upper()


def _is_title_shape(shape) -> bool:
    name = _shape_name(shape)
    if "SUBTITLE" in name:
        return False
    if "TITLE" in name:
        return True
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        ph = str(shape.placeholder_format.type).upper()
    except Exception:
        return False
    return "TITLE" in ph and "SUBTITLE" not in ph


def _is_subtitle_shape(shape) -> bool:
    name = _shape_name(shape)
    if "SUBTITLE" in name:
        return True
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        ph = str(shape.placeholder_format.type).upper()
    except Exception:
        return False
    return "SUBTITLE" in ph


def _is_content_shape(shape) -> bool:
    if _is_title_shape(shape) or _is_subtitle_shape(shape):
        return False
    name = _shape_name(shape)
    if any(token in name for token in ("BODY", "CONTENT", "TEXT", "OBJECT")):
        return True
    return bool(getattr(shape, "has_text_frame", False))


def _style_text_runs(text_frame, *, color: RGBColor, size_pt: float | None = None, bold: bool | None = None) -> None:
    for paragraph in text_frame.paragraphs:
        if size_pt is not None:
            paragraph.font.size = Pt(size_pt)
        if bold is not None:
            paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        for run in paragraph.runs:
            if size_pt is not None:
                run.font.size = Pt(size_pt)
            if bold is not None:
                run.font.bold = bold
            run.font.color.rgb = color


def _ensure_card_fill(shape, *, fill_color: RGBColor, border_color: RGBColor) -> None:
    if not hasattr(shape, "fill") or not hasattr(shape, "line"):
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.25)


def _apply_slide_theme(slide, theme: str, *, is_title_slide: bool, slide_width) -> None:
    palette = _THEME_COLORS.get((theme or "default").lower(), _THEME_COLORS["default"])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette["bg"]

    # Accent strip gives each slide a visible design motif even without custom template files.
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_width, Inches(0.16))
    strip.fill.solid()
    strip.fill.fore_color.rgb = palette["accent"]
    strip.line.fill.background()

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        if _is_title_shape(shape):
            _ensure_card_fill(shape, fill_color=palette["accent_soft"], border_color=palette["accent"])
            _style_text_runs(shape.text_frame, color=palette["title"], size_pt=34 if is_title_slide else 28, bold=True)
            continue

        if _is_subtitle_shape(shape):
            _ensure_card_fill(shape, fill_color=palette["card_bg"], border_color=palette["card_border"])
            _style_text_runs(shape.text_frame, color=palette["body"], size_pt=16, bold=False)
            continue

        if _is_content_shape(shape):
            _ensure_card_fill(shape, fill_color=palette["card_bg"], border_color=palette["card_border"])
            _style_text_runs(shape.text_frame, color=palette["body"], size_pt=14, bold=False)


def _write_lines(text_frame, lines: list[str]) -> None:
    text_frame.clear()
    if not lines:
        return
    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0


def _fill_title_slide(slide, slots: dict[str, str], title: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = slots.get("TITLE", title[:90])
    if len(slide.placeholders) > 1 and getattr(slide.placeholders[1], "has_text_frame", False):
        subtitle = (
            slots.get("SUBTITLE")
            or slots.get("BODY_1")
            or slots.get("TEXT_1")
            or "Generated by PowerPoint Agent"
        )
        slide.placeholders[1].text = subtitle


def _fill_content_slide(slide, slots: dict[str, str]) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = slots.get("TITLE", "Overview")

    content_shapes = [
        shape
        for shape in slide.placeholders
        if getattr(shape, "has_text_frame", False)
        and not (getattr(shape, "is_placeholder", False) and "TITLE" in str(getattr(shape, "name", "")).upper())
    ]

    bullet_lines = _normalize_lines(slots.get("BULLET_1", ""))
    body_lines = _normalize_lines(slots.get("BODY_1", ""))
    citation_lines = _normalize_lines(slots.get("CITATION", ""))
    citation_lines = [line if line.lower().startswith("source:") else f"Source: {line}" for line in citation_lines]
    primary_lines = bullet_lines + body_lines

    if not content_shapes:
        textbox = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.5))
        _write_lines(textbox.text_frame, primary_lines + citation_lines)
        return

    if len(content_shapes) >= 2:
        _write_lines(content_shapes[0].text_frame, primary_lines or body_lines or bullet_lines)
        _write_lines(content_shapes[1].text_frame, citation_lines)
    else:
        _write_lines(content_shapes[0].text_frame, primary_lines + citation_lines)


def build_scratch_pptx(
    *,
    slides_payload: list[dict],
    output_path: Path,
    title: str,
    theme: str = "default",
) -> None:
    prs = _presentation_for_theme(theme)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if not slides_payload:
        slide = prs.slides.add_slide(prs.slide_layouts[_pick_layout_index(prs, "section_break", 0)])
        _fill_title_slide(slide, {"TITLE": title[:90]}, title)
        _apply_slide_theme(slide, theme, is_title_slide=True, slide_width=prs.slide_width)
        prs.save(str(output_path))
        return

    for idx, row in enumerate(slides_payload):
        slots = {str(k).upper(): str(v) for k, v in (row.get("slots") or {}).items()}
        archetype = str(row.get("archetype") or "general")
        layout_index = _pick_layout_index(prs, archetype, idx)
        adjusted = _adjust_layout_for_content(prs, slots) if idx > 0 else None
        if adjusted is not None:
            layout_index = adjusted
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])

        if idx == 0:
            _fill_title_slide(slide, slots, title)
        else:
            _fill_content_slide(slide, slots)

        _apply_slide_theme(slide, theme, is_title_slide=idx == 0, slide_width=prs.slide_width)

    prs.save(str(output_path))
