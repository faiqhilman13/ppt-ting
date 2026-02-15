import re
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.slide_archetypes import infer_archetype

TOKEN_PATTERN = re.compile(r"\{\{\s*([A-Za-z0-9_:-]+)\s*\}\}")
NON_CONTENT_PLACEHOLDER_TAGS = ("DATE", "FOOTER", "SLIDE_NUMBER", "HEADER")
EMU_PER_INCH = 914400


def _placeholder_type_name(shape) -> str:
    try:
        return str(shape.placeholder_format.type).upper()
    except Exception:
        return str(getattr(shape, "name", "")).upper()


def _is_non_content_placeholder(shape) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    ph_type = _placeholder_type_name(shape)
    return any(tag in ph_type for tag in NON_CONTENT_PLACEHOLDER_TAGS)


def _extract_tokens(text: str) -> list[str]:
    if not text:
        return []
    return sorted({match.group(1).upper() for match in TOKEN_PATTERN.finditer(text)})


def _iter_shapes(shape_collection):
    for shape in shape_collection:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _shape_sort_key(shape) -> tuple[float, float, int]:
    left = float(getattr(shape, "left", 0))
    top = float(getattr(shape, "top", 0))
    shape_id = int(getattr(shape, "shape_id", -1))
    return (top, left, shape_id)


def _inches(raw_value: int | float | None) -> float:
    if not raw_value:
        return 0.0
    return round(float(raw_value) / EMU_PER_INCH, 2)


def _extract_font_size_pt(shape) -> float | None:
    if not getattr(shape, "has_text_frame", False):
        return None

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            size = run.font.size
            if size is not None:
                try:
                    return round(float(size.pt), 1)
                except Exception:
                    continue
    return None


def _shape_context(shape, text: str) -> dict:
    context = {
        "left_inches": _inches(getattr(shape, "left", 0)),
        "top_inches": _inches(getattr(shape, "top", 0)),
        "width_inches": _inches(getattr(shape, "width", 0)),
        "height_inches": _inches(getattr(shape, "height", 0)),
        "existing_text": (text or "")[:220],
    }
    font_size_pt = _extract_font_size_pt(shape)
    if font_size_pt:
        context["font_size_pt"] = font_size_pt
    return context


def _auto_slot_name(shape, counters: dict[str, int]) -> str | None:
    if not getattr(shape, "is_placeholder", False):
        return None

    if _is_non_content_placeholder(shape):
        return None

    ph_type = _placeholder_type_name(shape)

    if "TITLE" in ph_type and "SUBTITLE" not in ph_type:
        counters["TITLE"] += 1
        return "TITLE" if counters["TITLE"] == 1 else f"TITLE_{counters['TITLE']}"
    if "SUBTITLE" in ph_type:
        counters["SUBTITLE"] += 1
        return "SUBTITLE" if counters["SUBTITLE"] == 1 else f"SUBTITLE_{counters['SUBTITLE']}"

    counters["BODY"] += 1
    return f"BODY_{counters['BODY']}"


def parse_template_manifest(template_path: Path) -> dict:
    presentation = Presentation(str(template_path))

    slides: list[dict] = []
    all_slots: set[str] = set()

    for slide_index, slide in enumerate(presentation.slides):
        slot_counters: dict[str, int] = defaultdict(int)
        bindings: list[dict] = []

        for shape in sorted(_iter_shapes(slide.shapes), key=_shape_sort_key):
            shape_id = int(getattr(shape, "shape_id", -1))
            shape_name = str(getattr(shape, "name", f"shape-{shape_id}"))

            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text or ""

                tokens = _extract_tokens(text)
                for token in tokens:
                    bindings.append(
                        {
                            "slot": token,
                            "kind": "token_text",
                            "shape_id": shape_id,
                            "shape_name": shape_name,
                            **_shape_context(shape, text),
                        }
                    )
                    all_slots.add(token)

                # Auto-bind placeholders when explicit tokens are absent.
                if not tokens:
                    auto_slot = _auto_slot_name(shape, slot_counters)
                    if not auto_slot and text.strip() and not _is_non_content_placeholder(shape):
                        # Inventory-style fallback for non-placeholder text shapes.
                        slot_counters["TEXT"] += 1
                        auto_slot = f"TEXT_{slot_counters['TEXT']}"
                    if auto_slot:
                        bindings.append(
                            {
                                "slot": auto_slot,
                                "kind": "shape_text",
                                "shape_id": shape_id,
                                "shape_name": shape_name,
                                **_shape_context(shape, text),
                            }
                        )
                        all_slots.add(auto_slot)

            if getattr(shape, "has_table", False):
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        tokens = _extract_tokens(cell.text or "")
                        for token in tokens:
                            bindings.append(
                                {
                                    "slot": token,
                                    "kind": "token_table",
                                    "shape_id": shape_id,
                                    "shape_name": shape_name,
                                    "row": row_idx,
                                    "col": col_idx,
                                    **_shape_context(shape, cell.text or ""),
                                }
                            )
                            all_slots.add(token)
                        if not tokens and (cell.text or "").strip():
                            slot_counters["TABLE"] += 1
                            auto_slot = f"TABLE_{slot_counters['TABLE']}"
                            bindings.append(
                                {
                                    "slot": auto_slot,
                                    "kind": "table_cell_text",
                                    "shape_id": shape_id,
                                    "shape_name": shape_name,
                                    "row": row_idx,
                                    "col": col_idx,
                                    **_shape_context(shape, cell.text or ""),
                                }
                            )
                            all_slots.add(auto_slot)

        # De-duplicate exact bindings.
        unique_bindings = []
        seen = set()
        for binding in bindings:
            key = (
                binding["slot"],
                binding["kind"],
                binding["shape_id"],
                binding.get("row"),
                binding.get("col"),
            )
            if key in seen:
                continue
            seen.add(key)
            unique_bindings.append(binding)

        slide_slots = sorted({binding["slot"] for binding in unique_bindings})
        archetype = infer_archetype(slide_slots)

        slides.append(
            {
                "index": slide_index,
                "name": slide.slide_layout.name if slide.slide_layout else f"slide-{slide_index}",
                "slots": slide_slots,
                "archetype": archetype,
                "bindings": unique_bindings,
            }
        )

    return {
        "version": "template_fidelity_v1",
        "slide_count": len(slides),
        "slides": slides,
        "slot_types": sorted(all_slots),
    }


def _find_shape_by_id(slide, shape_id: int):
    for shape in _iter_shapes(slide.shapes):
        if int(getattr(shape, "shape_id", -1)) == int(shape_id):
            return shape
    return None


def extract_current_slot_values(pptx_path: Path, manifest: dict) -> list[dict]:
    presentation = Presentation(str(pptx_path))
    extracted: list[dict] = []

    for slide_spec in manifest.get("slides", []):
        idx = int(slide_spec.get("index", -1))
        if idx < 0 or idx >= len(presentation.slides):
            continue

        slide = presentation.slides[idx]
        slot_values: dict[str, str] = {}

        for binding in slide_spec.get("bindings", []):
            slot = str(binding.get("slot", "")).strip()
            if not slot:
                continue

            shape = _find_shape_by_id(slide, int(binding.get("shape_id", -1)))
            if not shape:
                continue

            kind = str(binding.get("kind", ""))
            value = ""
            if kind in {"shape_text", "token_text"} and getattr(shape, "has_text_frame", False):
                value = shape.text_frame.text or ""
            elif kind in {"table_cell_text", "token_table"} and getattr(shape, "has_table", False):
                row_idx = int(binding.get("row", -1))
                col_idx = int(binding.get("col", -1))
                if row_idx >= 0 and col_idx >= 0 and row_idx < len(shape.table.rows) and col_idx < len(shape.table.columns):
                    value = shape.table.cell(row_idx, col_idx).text or ""

            if not value.strip():
                continue

            if slot in slot_values and value != slot_values[slot]:
                slot_values[slot] = f"{slot_values[slot]}\n{value}"
            else:
                slot_values[slot] = value

        extracted.append({"template_slide_index": idx, "slots": slot_values})

    return extracted
