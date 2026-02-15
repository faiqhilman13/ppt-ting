from __future__ import annotations

import re
import shutil
from copy import deepcopy
from pathlib import Path

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

TOKEN_PATTERN = re.compile(r"\{\{\s*([A-Za-z0-9_:-]+)\s*\}\}")
NON_CONTENT_PLACEHOLDER_TAGS = ("DATE", "FOOTER", "SLIDE_NUMBER", "HEADER")

app = FastAPI(title="Template Fidelity PPTX Renderer")


class SlideInstruction(BaseModel):
    template_slide_index: int
    slots: dict[str, str] = Field(default_factory=dict)


class RenderRequest(BaseModel):
    deckId: str
    version: int
    slides: list[SlideInstruction]
    templateManifest: dict = Field(default_factory=dict)
    templatePath: str | None = None
    basePptxPath: str | None = None
    outputPath: str


@app.get("/health")
def health():
    return {"status": "ok"}


def _iter_shapes(shape_collection):
    for shape in shape_collection:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _replace_tokens(text: str, token_map: dict[str, str]) -> str:
    if not text:
        return text

    def repl(match):
        key = match.group(1).upper()
        return token_map.get(key, match.group(0))

    return TOKEN_PATTERN.sub(repl, text)


def _paragraph_text(paragraph) -> str:
    runs = list(paragraph.runs)
    if runs:
        return "".join(run.text or "" for run in runs)
    return paragraph.text or ""


def _write_paragraph_text_preserve_runs(paragraph, new_text: str) -> None:
    runs = list(paragraph.runs)
    if not runs:
        paragraph.text = new_text
        return

    remaining = new_text
    for idx, run in enumerate(runs):
        if idx == len(runs) - 1:
            run.text = remaining
            break

        current = run.text or ""
        take = min(len(current), len(remaining))
        run.text = remaining[:take]
        remaining = remaining[take:]


def _clear_paragraph_bullets(paragraph) -> None:
    p_pr = paragraph._p.find("{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
    if p_pr is None:
        return
    for child in list(p_pr):
        local_name = child.tag.split("}")[-1]
        if local_name.startswith("bu"):
            p_pr.remove(child)


def _write_text_frame_preserve_format(text_frame, new_text: str) -> None:
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = new_text
        return

    lines = str(new_text or "").replace("\r\n", "\n").split("\n")
    if not lines:
        lines = [""]

    # When more lines are provided than existing paragraphs, clone the last paragraph
    # so overflow lines keep paragraph/run-level formatting instead of collapsing with soft breaks.
    while len(paragraphs) < len(lines):
        template_paragraph = paragraphs[-1]
        text_frame._txBody.append(deepcopy(template_paragraph._p))
        paragraphs = list(text_frame.paragraphs)

    for idx, paragraph in enumerate(paragraphs):
        if idx >= len(lines):
            break
        _write_paragraph_text_preserve_runs(paragraph, lines[idx])

    # Remove stale trailing paragraphs when the new content has fewer lines.
    # This avoids empty bullets/extra vertical spacing from template leftovers.
    if len(lines) < len(paragraphs):
        for idx in range(len(paragraphs) - 1, len(lines) - 1, -1):
            if idx <= 0:
                break
            text_frame._txBody.remove(paragraphs[idx]._p)

    # When content is intentionally blank, clear bullet marker properties from
    # the remaining paragraph to avoid orphan bullets in templates.
    if len(lines) == 1 and not lines[0].strip():
        refreshed = list(text_frame.paragraphs)
        if refreshed:
            _clear_paragraph_bullets(refreshed[0])


def _placeholder_type_name(shape) -> str:
    try:
        return str(shape.placeholder_format.type).upper()
    except Exception:
        return str(getattr(shape, "name", "")).upper()


def _is_non_content_placeholder(shape) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    ph_type = _placeholder_type_name(shape)
    return any(tag in ph_type for tag in NON_CONTENT_PLACEHOLDER_TAGS)


def _apply_token_replacements(slide, token_map: dict[str, str]) -> None:
    for shape in _iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            original_text = shape.text_frame.text or ""
            updated_text = _replace_tokens(original_text, token_map)
            if updated_text != original_text:
                _write_text_frame_preserve_format(shape.text_frame, updated_text)

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if not getattr(cell, "text_frame", None):
                        continue
                    original_text = cell.text_frame.text or ""
                    updated_text = _replace_tokens(original_text, token_map)
                    if updated_text != original_text:
                        _write_text_frame_preserve_format(cell.text_frame, updated_text)


def _find_shape_by_id(slide, shape_id: int):
    for shape in _iter_shapes(slide.shapes):
        if int(getattr(shape, "shape_id", -1)) == shape_id:
            return shape
    return None


def _apply_shape_bindings(slide, slide_manifest: dict, slot_values: dict[str, str]) -> None:
    for binding in slide_manifest.get("bindings", []):
        slot_name = str(binding.get("slot", "")).upper()
        if slot_name not in slot_values:
            continue

        shape_id = int(binding.get("shape_id", -1))
        shape = _find_shape_by_id(slide, shape_id)
        if not shape:
            continue

        kind = str(binding.get("kind", ""))
        if kind == "shape_text":
            if not getattr(shape, "has_text_frame", False):
                continue
            if _is_non_content_placeholder(shape):
                continue
            _write_text_frame_preserve_format(shape.text_frame, slot_values[slot_name])
            continue

        if kind == "table_cell_text":
            if not getattr(shape, "has_table", False):
                continue
            row_idx = int(binding.get("row", -1))
            col_idx = int(binding.get("col", -1))
            if row_idx < 0 or col_idx < 0:
                continue
            if row_idx >= len(shape.table.rows):
                continue
            if col_idx >= len(shape.table.columns):
                continue
            cell = shape.table.cell(row_idx, col_idx)
            if not getattr(cell, "text_frame", None):
                continue
            _write_text_frame_preserve_format(cell.text_frame, slot_values[slot_name])


def _prune_unselected_slides(prs: Presentation, selected_indices: set[int]) -> None:
    # Keep selected slides only to avoid leaking untouched template content.
    for idx in range(len(prs.slides) - 1, -1, -1):
        if idx in selected_indices:
            continue
        slide_id = prs.slides._sldIdLst[idx].rId
        prs.part.drop_rel(slide_id)
        del prs.slides._sldIdLst[idx]


def _render(req: RenderRequest) -> None:
    source_path = Path(req.basePptxPath or req.templatePath or "")
    if not source_path.exists():
        raise FileNotFoundError(f"Source PPTX does not exist: {source_path}")

    output_path = Path(req.outputPath)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(source_path, output_path)

    prs = Presentation(str(output_path))
    manifest_by_index = {int(slide.get("index", i)): slide for i, slide in enumerate(req.templateManifest.get("slides", []))}

    for instruction in req.slides:
        idx = instruction.template_slide_index
        if idx < 0 or idx >= len(prs.slides):
            continue

        slide = prs.slides[idx]
        slot_values = {k.upper(): str(v) for k, v in instruction.slots.items()}
        _apply_token_replacements(slide, slot_values)
        _apply_shape_bindings(slide, manifest_by_index.get(idx, {}), slot_values)

    selected_indices = {int(item.template_slide_index) for item in req.slides}
    _prune_unselected_slides(prs, selected_indices)

    prs.save(str(output_path))


@app.post("/render")
def render(req: RenderRequest):
    try:
        _render(req)
        return {"status": "ok", "outputPath": req.outputPath}
    except FileNotFoundError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Render failed: {exc}") from exc
